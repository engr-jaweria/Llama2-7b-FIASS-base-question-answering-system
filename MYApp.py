import os
import tempfile
import numpy as np
import faiss
import requests
import pandas as pd
from sentence_transformers import SentenceTransformer
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
import streamlit as st
import fitz  # PyMuPDF for PDF handling
from docx import Document as DocxDocument
from pptx import Presentation
from dotenv import load_dotenv
from langchain.schema import Document
from langchain.docstore import InMemoryDocstore
from langchain.vectorstores import FAISS as LangchainFAISS
from langchain.docstore.document import Document as LangchainDocument
import replicate

# Load environment variables
load_dotenv()

# Set environment variables
REPLICATE_API_TOKEN = os.getenv("REPLICATE_API_TOKEN")
HUGGINGFACE_API_TOKEN = os.getenv("HUGGINGFACE_API_TOKEN")
if not REPLICATE_API_TOKEN or not HUGGINGFACE_API_TOKEN:
    raise ValueError("API tokens for Replicate and Hugging Face are not set in environment variables.")

# Initialize models
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
embeddings_model = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

# Define utility functions
def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file."""
    text = ""
    with fitz.open(pdf_file) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(docx_file):
    """Extract text from a DOCX file."""
    doc = DocxDocument(docx_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    """Extract text from a PPTX file."""
    presentation = Presentation(pptx_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(xlsx_file):
    """Extract text from an XLSX file."""
    df = pd.read_excel(xlsx_file)
    return df.to_string()

def extract_text_from_url(url):
    """Fetch text from a URL."""
    try:
        response = requests.get(url)
        response.raise_for_status()
        return response.text
    except requests.exceptions.RequestException as e:
        st.error(f"Error fetching URL content: {e}")
        return ""

def validate_total_file_size(files, max_size_mb=20):
    """Validate the total size of uploaded files."""
    total_size = sum(file.size for file in files) / (1024 * 1024)  # Convert bytes to MB
    if total_size > max_size_mb:
        st.error(f"Total file size exceeds {max_size_mb} MB. Please reduce the total file size.")
        return False
    return True

# Hugging Face Llama2 API integration
def llama2_generate_huggingface(prompt, top_p=1, temperature=0.75, max_new_tokens=800):
    """Generate text using Hugging Face Llama2 API."""
    url = "https://api-inference.huggingface.co/models/meta-llama/Llama-2-7b-hf"
    headers = {
        "Authorization": f"Bearer {HUGGINGFACE_API_TOKEN}"
    }
    payload = {
        "inputs": prompt,
        "parameters": {
            "top_p": top_p,
            "temperature": temperature,
            "max_new_tokens": max_new_tokens
        }
    }
    response = requests.post(url, headers=headers, json=payload)
    
    if response.status_code == 200:
        result = response.json()
        return result[0]['generated_text']  # The model response will be here
    else:
        st.error(f"Error with Hugging Face API: {response.text}")
        return "Error generating answer."

# Replicate Llama2 API integration
def llama2_generate_replicate(prompt, top_p=1, temperature=0.75, max_new_tokens=800):
    """Generate text using Replicate Llama2 API."""
    model = replicate.models.get("replicate/llama2-7b")
    output = model.predict(
        prompt=prompt,
        top_p=top_p,
        temperature=temperature,
        max_new_tokens=max_new_tokens
    )
    return output[0]  # Assuming the model returns a list of responses

# Main Streamlit application
def main():
    st.title("Llama2-7B Q&A System")

    # Sidebar inputs
    st.sidebar.header("Upload or Provide Input")
    uploaded_files = st.sidebar.file_uploader("Upload files", type=["txt", "pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)
    web_links = st.sidebar.text_area("Or provide web links (comma-separated):", "").split(',')

    # Allow user to limit number of uploaded files
    chunk_size = st.sidebar.slider("Set Chunk Size", 500, 2000, 1000)

    documents = []

    # Validate and process uploaded files
    if uploaded_files:
        if not validate_total_file_size(uploaded_files):
            return

        temp_dir = tempfile.TemporaryDirectory()

        for uploaded_file in uploaded_files:  # Process all uploaded files
            file_path = os.path.join(temp_dir.name, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.read())

            if uploaded_file.name.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    documents.append(f.read())
            elif uploaded_file.name.endswith(".pdf"):
                documents.append(extract_text_from_pdf(file_path))
            elif uploaded_file.name.endswith(".docx"):
                documents.append(extract_text_from_docx(file_path))
            elif uploaded_file.name.endswith(".pptx"):
                documents.append(extract_text_from_pptx(file_path))
            elif uploaded_file.name.endswith(".xlsx"):
                documents.append(extract_text_from_xlsx(file_path))

        temp_dir.cleanup()

    # Process web links
    for web_link in web_links:
        if web_link.strip():  # Skip empty links
            documents.append(extract_text_from_url(web_link.strip()))

    if documents:
        # Split documents into chunks without overlap
        chunks = [doc[i:i+chunk_size] for doc in documents for i in range(0, len(doc), chunk_size)]
        docs = [Document(page_content=chunk, metadata={"source": "uploaded file"}) for chunk in chunks]

        # Generate embeddings
        embeddings = embeddings_model.embed_documents(chunks)

        # Convert embeddings to NumPy array (shape: num_embeddings x embedding_dim)
        embeddings_array = np.array(embeddings)
        dimension = embeddings_array.shape[1]
        faiss_index = faiss.IndexFlatL2(dimension)
        faiss_index.add(embeddings_array.astype(np.float32))

        # Create a docstore
        docstore = InMemoryDocstore({i: docs[i] for i in range(len(docs))})

        # Prepare LangChain documents
        vector_store = LangchainFAISS.from_documents(docs, embeddings_model)

        # Set up retriever
        retriever = vector_store.as_retriever(search_kwargs={"k": 5})

        # Chat interface
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []

        st.header("Ask a Question")
        question = st.text_input("Your Question:")
        format_choice = st.selectbox("Answer Format:", ["Default", "Bullet Points", "Summary", "Specific Length"])
        word_limit = st.number_input("Word/Character Limit:", min_value=10, max_value=500, value=100) if format_choice == "Specific Length" else None

        model_choice = st.selectbox("Choose model:", ["Hugging Face", "Replicate"])

        if st.button("Get Answer") and question:
            context_docs = retriever.get_relevant_documents(question)
            context_docs.sort(key=lambda x: x.metadata.get("relevance", 0), reverse=True)  # Sort by relevance
            context = " ".join([doc.page_content for doc in context_docs])
            prompt = f"Answer the question based on the context:\n{context}\n\nQuestion: {question}"
        
            # Modify the prompt to suit different answer formats
            if format_choice == "Bullet Points":
                prompt += "\nPlease provide the answer in bullet points."
            elif format_choice == "Summary":
                prompt += "\nPlease provide a concise summary."
            elif format_choice == "Specific Length":
                prompt += f"\nPlease limit the answer to {word_limit} words."
        
            # Call the appropriate model for generation
            if model_choice == "Hugging Face":
                answer = llama2_generate_huggingface(prompt)
            elif model_choice == "Replicate":
                answer = llama2_generate_replicate(prompt)
            else:
                answer = "Please select a model."
        
            # Format the answer according to the chosen style
            if format_choice == "Bullet Points":
                answer = "\n- " + answer.replace(". ", "\n- ")
            elif format_choice == "Summary":
                answer = "\nSummary: " + answer
            elif format_choice == "Specific Length":
                answer = answer[:word_limit]
        
            # Save chat history
            st.session_state.chat_history.append({"question": question, "answer": answer})
        
            # Display question and answer
            st.markdown("## **Generated Response:**")
            st.markdown(f"### **{question}**")
            st.write(answer)
        
        # Display chat history (Only if the user scrolls to it explicitly)
        if st.session_state.chat_history:
            with st.expander("Chat History"):
                st.header("Chat History")
                for entry in st.session_state.chat_history:
                    st.write(f"**Q:** {entry['question']}\n**A:** {entry['answer']}\n")


    else:
        st.write("Upload documents or provide a web link to begin.")

if __name__ == "__main__":
    main()
